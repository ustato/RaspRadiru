# -*- coding: utf-8 -*-

from pptx import Presentation
import datetime


def make_testPPTX(test,lines) :
    # 日付取得
    main_title = "ラジオ英会話問題"
    today = datetime.date.today()

    # 初期化とスタートページ作成
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = main_title
    subtitle.text = str(today)


    # 問題と正解を作るためにリストを変形
    test=list(test)
    test*=2
    test.sort()
    for (i,j) in enumerate(test) :
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        if i%2 :
            title.text = lines[j-2]
            subtitle.text = lines[j-1]
        else :
            title.text = ""
            subtitle.text = lines[j-1]

    # 保存
    prs.save('test.pptx')
